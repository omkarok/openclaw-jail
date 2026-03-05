import subprocess
subprocess.run(["pip3","install","python-pptx","-q","--break-system-packages"],check=False)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
slides_data=[{"title": "Slide 1 Title", "headline": "Key message 1", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 2 Title", "headline": "Key message 2", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 3 Title", "headline": "Key message 3", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 4 Title", "headline": "Key message 4", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 5 Title", "headline": "Key message 5", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 6 Title", "headline": "Key message 6", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 7 Title", "headline": "Key message 7", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 8 Title", "headline": "Key message 8", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 9 Title", "headline": "Key message 9", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 10 Title", "headline": "Key message 10", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 11 Title", "headline": "Key message 11", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 12 Title", "headline": "Key message 12", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 13 Title", "headline": "Key message 13", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 14 Title", "headline": "Key message 14", "bullets": ["Point A", "Point B", "Point C"]}, {"title": "Slide 15 Title", "headline": "Key message 15", "bullets": ["Point A", "Point B", "Point C"]}]
prs=Presentation()
prs.slide_width=Inches(13.33)
prs.slide_height=Inches(7.5)
for i,s in enumerate(slides_data, start=1):
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    bg=slide.background.fill
    bg.solid(); bg.fore_color.rgb=RGBColor(15,15,25)
    tb=slide.shapes.add_textbox(Inches(0.6),Inches(0.3),Inches(12.1),Inches(0.9)).text_frame
    p=tb.paragraphs[0]; r=p.add_run(); r.text=s.get('title',''); r.font.size=Pt(36); r.font.bold=True; r.font.color.rgb=RGBColor(255,255,255)
    hb=slide.shapes.add_textbox(Inches(0.65),Inches(1.25),Inches(12),Inches(0.7)).text_frame
    p=hb.paragraphs[0]; r=p.add_run(); r.text=s.get('headline',''); r.font.size=Pt(20); r.font.italic=True; r.font.color.rgb=RGBColor(99,179,237)
    bb=slide.shapes.add_textbox(Inches(0.9),Inches(2.0),Inches(11.8),Inches(4.8)).text_frame
    bb.word_wrap=True
    first=True
    for b in s.get('bullets',[]) or ['']:
        p=bb.paragraphs[0] if first else bb.add_paragraph(); first=False
        p.text='• '+b if b else ''
        p.level=0
        if p.runs:
            p.runs[0].font.size=Pt(18); p.runs[0].font.color.rgb=RGBColor(220,220,220)
    nb=slide.shapes.add_textbox(Inches(12.2),Inches(7.0),Inches(1.0),Inches(0.3)).text_frame
    p=nb.paragraphs[0]; r=p.add_run(); r.text=str(i); r.font.size=Pt(12); r.font.color.rgb=RGBColor(100,100,120)
prs.save('/home/node/workspace/results/ma-explainer.pptx')
print('saved')
