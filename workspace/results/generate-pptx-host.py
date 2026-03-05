from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG    = RGBColor(15,  15,  25)
WHITE = RGBColor(255, 255, 255)
ACCENT= RGBColor(99,  179, 237)
BODY  = RGBColor(220, 220, 220)
DIM   = RGBColor(100, 100, 120)
GREEN = RGBColor(74,  222, 128)

slides_data = [
  {
    "title": "The AI That Works While You Sleep",
    "headline": "While you were asleep, your AI completed 8 tasks, upgraded its own mandate, and messaged you the results.",
    "bullets": [
      "OpenClaw: cross-channel AI runtime and security-hardened gateway",
      "Sherbyte: personal AI bridge — always listening, always delivering",
      "Background Worker: autonomous agent that executes while you're offline",
      "0 task failures across the entire production run",
    ],
    "notes": "Open with the hook — this is running in production. The numbers are real.",
  },
  {
    "title": "The Problem",
    "headline": "Every AI assistant you've used resets when you close the tab — that's not intelligence, that's a search bar with a personality.",
    "bullets": [
      "Reactive-only: AI waits for a message; it never initiates",
      "Stateless by default: close the session, lose the thread forever",
      "No execution layer: AI can plan but cannot run, track, or retry",
      "No proactive delivery: results sit in a chat window until you check",
    ],
    "notes": "Reactive and stateless are the two core failure modes. The missing piece is runtime infrastructure, not model quality.",
  },
  {
    "title": "Our Answer",
    "headline": "Three layers: a runtime, a personal bridge, and a headless worker — each doing exactly one job, autonomously.",
    "bullets": [
      "Layer 1 — OpenClaw: multi-channel gateway, tool routing, security enforcement",
      "Layer 2 — Sherbyte: real-time agent, heartbeat loop, notification delivery",
      "Layer 3 — Background Worker: cron-driven executor, reads MANDATE.md, drains queue",
      "Layers communicate via flat JSON files — auditable, crash-safe, zero coupling",
    ],
    "notes": "Key insight: each layer is independently replaceable. The communication contract is just JSON files.",
  },
  {
    "title": "Architecture Diagram",
    "headline": "One gateway, three channels, two agents, five JSON files — the entire system fits on one page.",
    "bullets": [
      "OpenClaw gateway: ws://127.0.0.1:18789 — localhost-only, hardened Docker jail",
      "Sherbyte: named agent, 30-min heartbeat, notification + escalation delivery",
      "Background Worker: headless, cron 08:00 IST, MANDATE.md v1.4",
      "State: queue.json, notifications.json, escalations.json, MANDATE.md, results/",
    ],
    "notes": "Walk left to right: WhatsApp input, OpenClaw routes, Sherbyte queues, Worker executes, Sherbyte delivers.",
    "code": (
      "  USER CHANNELS\n"
      "  WhatsApp   Discord   Web/API\n"
      "      down       down      down\n"
      "  OPENCLAW GATEWAY :18789\n"
      "  0 critical  0 warn  uid=1000\n"
      "         down\n"
      "     left    right\n"
      "  SHERBYTE     BACKGROUND WORKER\n"
      "  30-min       Cron 08:00 IST\n"
      "  heartbeat    MANDATE.md v1.4\n"
      "       down         down\n"
      "  queue.json   notifications.json\n"
      "  escalations.json   results/\n"
      "              down\n"
      "       WhatsApp DM to Owner"
    ),
  },
  {
    "title": "OpenClaw — The Runtime Fabric",
    "headline": "OpenClaw is the surface the agents live on — it handles channels, security, and tool routing so agents don't have to.",
    "bullets": [
      "Binds at ws://127.0.0.1:18789 — localhost-only; no public exposure by default",
      "Channels: WhatsApp (Baileys), Discord, Web — each with independent policy",
      "dmPolicy: 'allowlist' + allowFrom per number — zero-trust DM access",
      "Security audit: 0 critical, 0 warn — read-only root FS, uid=1000, caps dropped",
    ],
    "notes": "Default dmPolicy allowlist means the bot responds to no one until you explicitly whitelist their number.",
  },
  {
    "title": "Sherbyte — Personal AI Bridge",
    "headline": "Sherbyte is the human-facing half: it listens in real time, escalates failures, and delivers results before you ask.",
    "bullets": [
      "Named agent in OpenClaw — persistent session, no context resets",
      "30-minute heartbeat: checks notifications.json + escalations.json, delivers, sleeps",
      "Parses !task commands from WhatsApp groups — writes structured entries to queue.json",
      "Delivers task results as WhatsApp DMs — proactively, no polling required",
    ],
    "notes": "The 30-minute heartbeat is what makes the system feel alive without burning continuous API quota.",
  },
  {
    "title": "Background Worker — Autonomous Executor",
    "headline": "No browser, no session, no human in the loop — the worker reads MANDATE.md and executes until the queue is empty.",
    "bullets": [
      "Headless agent: launched by cron at 08:00 IST daily, no UI",
      "Reads MANDATE.md (v1.4) on startup — constitution, not just a config file",
      "Drains queue.json in priority order: urgent, high, normal, low",
      "Crash recovery: any task in_progress > 30 min is reset and retried automatically",
    ],
    "notes": "MANDATE.md is human-readable plain English. The 30-min stale timeout means the system self-heals from crashes.",
    "code": (
      "[08:00:01 IST] Background Worker v1.4 starting...\n"
      "[08:00:02 IST] Reading MANDATE.md... OK (v1.4)\n"
      "[08:00:02 IST] Loading queue.json... 3 tasks pending\n"
      "[08:00:03 IST] Executing t009 [priority: urgent]...\n"
      "[08:00:47 IST] t009 complete. Writing result.\n"
      "[08:00:48 IST] Delivering notifications... 1 queued."
    ),
  },
  {
    "title": "Task Queue + DAG",
    "headline": "queue.json is not a list — it's a dependency graph with priorities, backoff, and self-healing crash recovery.",
    "bullets": [
      "Schema v2: each task has id, priority, depends_on[], status, attempts, run_after",
      "DAG via depends_on: tasks block until all dependencies reach status done",
      "Priority sort: urgent > high > normal > low — re-evaluated every drain cycle",
      "Exponential backoff: attempt 1 to 1min, 2 to 2min, 3 to 4min (max 1 hour)",
    ],
    "notes": "The backoff formula prevents a flaky API from hammering the worker into a retry storm.",
    "code": (
      '{\n'
      '  "id": "t007",\n'
      '  "title": "failure spike detection",\n'
      '  "priority": "high",\n'
      '  "status": "pending",\n'
      '  "depends_on": ["t005", "t006"],\n'
      '  "attempts": 0,\n'
      '  "max_attempts": 3,\n'
      '  "run_after": null\n'
      '}'
    ),
  },
  {
    "title": "The Self-Improvement Sprint",
    "headline": "The worker ran tasks t004 through t008, shipped four production features, and upgraded its own mandate three versions — overnight, alone.",
    "bullets": [
      "t004: Notifications system — worker to notifications.json to Sherbyte to WhatsApp DM",
      "t005: Per-attempt receipts — every task attempt logged with timestamp and result",
      "t006: Weekly digest — auto-generated summary of the week's task history",
      "t007 to t008: Failure spike detection + summary — MANDATE upgraded v1.1 to v1.4",
    ],
    "notes": "The worker identified gaps, created tasks to fill them, executed them, then updated its own MANDATE. No human wrote the upgrade.",
  },
  {
    "title": "!task — AI from Any WhatsApp Group",
    "headline": "Type !task research: X in a WhatsApp group, walk away — the result arrives in your DM before you remember you asked.",
    "bullets": [
      "User types !task research: X in any allowlisted WhatsApp group",
      "Sherbyte parses the prefix, extracts task type, writes structured entry to queue.json",
      "Background Worker picks it up next cycle, executes, writes result to results/",
      "Next heartbeat: Sherbyte finds sent:false, delivers DM, marks sent:true",
    ],
    "notes": "The !task prefix is the entire API surface for the end user. Full round-trip: approx 8 minutes, zero manual steps.",
    "code": (
      "[WhatsApp Group]    [Sherbyte]      [Worker]     [WhatsApp DM]\n"
      "!task research:  -> Parses prefix, -> Executes, -> Research done:\n"
      "quantum X           writes queue     writes result  results/t011\n"
      "\n"
      "          approx 8 minutes end-to-end, zero manual steps"
    ),
  },
  {
    "title": "The Notification Loop",
    "headline": "notifications.json is the delivery contract — a dead-simple guarantee with no message broker required.",
    "bullets": [
      'Worker writes: {"id": "n001", "message": "...", "sent": false}',
      "Sherbyte heartbeat reads all entries where sent == false every 30 minutes",
      "Delivers each notification to owner's WhatsApp DM via OpenClaw",
      "Marks sent: true — idempotent, crash-safe, no duplicate deliveries",
    ],
    "notes": "If Sherbyte crashes mid-delivery, notification stays sent:false and is delivered on next heartbeat. No Kafka, no Redis.",
    "code": (
      'Before:  {"id": "n005", "sent": false}\n'
      '         -> Sherbyte heartbeat fires\n'
      'After:   {"id": "n005", "sent": true,\n'
      '          "sent_at": "2026-03-04T09:50:52Z"}\n'
      "\n"
      "Delivery window: max 30 minutes."
    ),
  },
  {
    "title": "Production Results",
    "headline": "This is not a benchmark — these numbers are from a live system running in a Docker jail on a single laptop.",
    "bullets": [
      "8+ tasks completed autonomously — 0 failures, 0 manual retries, 0 human interventions",
      "5 WhatsApp notifications delivered — all marked sent:true, none duplicated",
      "MANDATE.md self-upgraded 3 times: v1.1 to v1.2 to v1.3 to v1.4 without human authorship",
      "Security posture held: 0 critical, 0 warn on every openclaw security audit run",
    ],
    "notes": "The MANDATE self-upgrade is the most surprising stat: the system wrote its own operating procedures as a byproduct of executing tasks.",
  },
  {
    "title": "What This Enables",
    "headline": "When AI has a runtime, a task queue, and a delivery layer, 'assistant' becomes the wrong word — you have an operator.",
    "bullets": [
      "Proactive AI: system initiates contact when something happens, not when you ask",
      "Auditable execution: every task has a JSON record — id, priority, attempts, result",
      "Ops infrastructure: the DAG queue pattern works for any autonomous agent, any model",
      "Compounding autonomy: each self-improvement sprint expands what the next can do",
    ],
    "notes": "The DAG queue is model-agnostic: swap Sherbyte for any LLM-backed agent and the orchestration layer stays identical.",
  },
  {
    "title": "How to Build This",
    "headline": "Three components, five JSON files, one cron job — you can wire this pattern in a weekend.",
    "bullets": [
      "Step 1 — Runtime: deploy OpenClaw, configure gateway.bind, add your channel",
      "Step 2 — Agent: define MANDATE.md with goals and constraints; connect your LLM",
      "Step 3 — Queue: implement queue.json schema v2 with depends_on, priority, run_after",
      "Glue: notifications.json (sent:false/true) + escalations.json (dedup_key)",
    ],
    "notes": "The hardest part is deciding what your agent's mandate actually is. Everything else is JSON and a cron job.",
  },
  {
    "title": "What's Next + Q&A",
    "headline": "The system already works — what it does next depends on what you put in MANDATE.md.",
    "bullets": [
      "Multi-worker parallelism: shard queue.json by priority, run concurrent workers",
      "Agent-to-agent delegation: Sherbyte spawns sub-tasks for Background Worker",
      "External triggers: GitHub webhooks, calendar events — anything that writes queue.json",
      "Open question: when should the agent refuse a MANDATE update it wrote itself?",
    ],
    "notes": "End on the open question — it's the most intellectually honest thing about autonomous systems and builds credibility.",
  },
]


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_tb(slide, text, l, t, w, h, color=WHITE, size=18,
           bold=False, italic=False, align=PP_ALIGN.LEFT, mono=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.color.rgb = color
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    if mono:
        r.font.name = "Courier New"
    return tb


prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
M = Inches(0.5)
W = prs.slide_width - M * 2

for i, d in enumerate(slides_data, 1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    has_code = "code" in d

    # Slide number
    add_tb(slide, f"{i}/15",
           prs.slide_width - Inches(1.1), Inches(0.1), Inches(0.9), Inches(0.35),
           DIM, 11, align=PP_ALIGN.RIGHT)

    # Title
    add_tb(slide, d["title"], M, Inches(0.25), W, Inches(0.8), WHITE, 30, bold=True)

    # Headline
    add_tb(slide, d["headline"], M, Inches(1.15), W, Inches(0.65), ACCENT, 15, italic=True)

    # Bullets
    bullet_h = Inches(2.2) if has_code else Inches(4.4)
    tb = slide.shapes.add_textbox(M, Inches(1.9), W, bullet_h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in d["bullets"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = "->  " + b
        r.font.color.rgb = BODY
        r.font.size = Pt(16)

    # Code block
    if has_code:
        tb2 = slide.shapes.add_textbox(M, Inches(4.2), W, Inches(2.85))
        tf2 = tb2.text_frame
        tf2.word_wrap = False
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = d["code"]
        r2.font.color.rgb = GREEN
        r2.font.size = Pt(10)
        r2.font.name = "Courier New"

    # Speaker notes
    if d.get("notes"):
        slide.notes_slide.notes_text_frame.text = d["notes"]

OUT = r"C:\Users\mentoria\openclaw-jail\workspace\results\ma-explainer.pptx"
prs.save(OUT)
print(f"Saved {len(slides_data)} slides -> {OUT}")
